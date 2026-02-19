"""Main generator: JSON spec to PPTX."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .html_parser import render_html_to_text_frame
from .schema import validate_spec
from .styles import parse_css_color, parse_css_font_size, parse_dimension

_LAYOUT_NAME_TO_INDEX = {
    "title": 0,
    "title_and_content": 1,
    "section_header": 2,
    "blank": 6,
}


def generate_pptx(spec: dict[str, Any], output_path: str | Path) -> Path:
    """Generate a PPTX file from a validated JSON-compatible dict."""
    spec = validate_spec(spec)

    presentation = Presentation()
    presentation_data = spec["presentation"]

    if "width" in presentation_data:
        presentation.slide_width = parse_dimension(presentation_data["width"])
    if "height" in presentation_data:
        presentation.slide_height = parse_dimension(presentation_data["height"])

    for slide_spec in spec["slides"]:
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[_resolve_layout_index(slide_spec.get("layout", "blank"))]
        )
        for element in slide_spec["elements"]:
            _render_element(slide, element, presentation_data.get("defaults") or {})

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return output_path


def generate_pptx_from_file(input_path: str | Path, output_path: str | Path) -> Path:
    """Load JSON from file and generate PPTX."""
    input_path = Path(input_path)
    with input_path.open("r", encoding="utf-8") as handle:
        spec = json.load(handle)
    return generate_pptx(spec, output_path)


def _resolve_layout_index(layout: str | int) -> int:
    if isinstance(layout, int):
        return layout
    return _LAYOUT_NAME_TO_INDEX[layout]


def _render_element(slide, element: dict[str, Any], presentation_defaults: dict[str, Any]) -> None:
    position = element["position"]
    textbox = slide.shapes.add_textbox(
        parse_dimension(position["left"]),
        parse_dimension(position["top"]),
        parse_dimension(position["width"]),
        parse_dimension(position["height"]),
    )
    if isinstance(element.get("name"), str) and element["name"].strip():
        textbox.name = element["name"].strip()

    content = element.get("content")
    if not content:
        return

    text_frame = textbox.text_frame
    element_style = element.get("style") or {}
    merged_style = {**presentation_defaults, **element_style}
    _apply_text_frame_style(text_frame, merged_style)
    render_html_to_text_frame(text_frame, content, _run_style_defaults(merged_style))
    _apply_paragraph_style(text_frame, merged_style)


def _run_style_defaults(style: dict[str, Any]) -> dict[str, Any]:
    defaults: dict[str, Any] = {}

    font_name = style.get("font_name")
    if isinstance(font_name, str) and font_name.strip():
        defaults["font_name"] = font_name.strip()

    font_size = style.get("font_size")
    if isinstance(font_size, str) and font_size.strip():
        try:
            defaults["font_size"] = parse_css_font_size(font_size)
        except ValueError:
            pass

    font_color = style.get("font_color")
    if isinstance(font_color, str) and font_color.strip():
        try:
            defaults["font_color"] = parse_css_color(font_color)
        except ValueError:
            pass

    return defaults


def _apply_text_frame_style(text_frame, style: dict[str, Any]) -> None:
    word_wrap = style.get("word_wrap")
    if isinstance(word_wrap, bool):
        text_frame.word_wrap = word_wrap

    vertical_anchor = style.get("vertical_anchor")
    if isinstance(vertical_anchor, str):
        normalized = vertical_anchor.strip().lower()
        mapping = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        if normalized in mapping:
            text_frame.vertical_anchor = mapping[normalized]


def _apply_paragraph_style(text_frame, style: dict[str, Any]) -> None:
    alignment = style.get("alignment")
    if not isinstance(alignment, str):
        return

    normalized = alignment.strip().lower()
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    if normalized not in mapping:
        return
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = mapping[normalized]
