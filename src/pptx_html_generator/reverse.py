"""Reverse conversion: PPTX shape text to HTML."""

from __future__ import annotations

import html
import re
from pathlib import Path
from typing import Any

from pptx import Presentation

_UL_PREFIXES = ("\u2022 ", "\u2013 ", "\u25e6 ")
_OL_PREFIX_RE = re.compile(r"^\d+\.\s+")


def list_slide_elements(pptx_path: str | Path, slide_number: int) -> list[dict[str, Any]]:
    """List shape metadata for a slide using Selection Pane names."""
    slide = _get_slide(pptx_path, slide_number)
    items: list[dict[str, Any]] = []
    for shape in slide.shapes:
        has_text = bool(getattr(shape, "has_text_frame", False))
        preview = ""
        if has_text:
            preview = (shape.text_frame.text or "").strip().replace("\n", " ")
            preview = preview[:80]
        items.append(
            {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "has_text": has_text,
                "text_preview": preview,
            }
        )
    return items


def extract_shape_html(
    pptx_path: str | Path,
    slide_number: int,
    shape_name: str | None = None,
    shape_id: int | None = None,
) -> str:
    """Extract HTML from a shape identified by Selection Pane name or shape id."""
    if not shape_name and shape_id is None:
        raise ValueError("Either shape_name or shape_id is required.")

    slide = _get_slide(pptx_path, slide_number)
    shape = _select_shape(slide, shape_name=shape_name, shape_id=shape_id)
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"Shape '{shape.name}' has no text frame.")
    return _text_frame_to_html(shape.text_frame)


def _get_slide(pptx_path: str | Path, slide_number: int):
    presentation = Presentation(str(pptx_path))
    if slide_number < 1 or slide_number > len(presentation.slides):
        raise ValueError(f"slide_number out of range: {slide_number}")
    return presentation.slides[slide_number - 1]


def _select_shape(slide, shape_name: str | None, shape_id: int | None):
    for shape in slide.shapes:
        if shape_name is not None and shape.name == shape_name:
            return shape
        if shape_id is not None and shape.shape_id == shape_id:
            return shape
    if shape_name is not None:
        raise ValueError(f"Shape not found by name: {shape_name}")
    raise ValueError(f"Shape not found by id: {shape_id}")


def _text_frame_to_html(text_frame) -> str:
    blocks: list[dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        paragraph_html = _paragraph_html(paragraph)
        blocks.append(_classify_paragraph(paragraph, paragraph_html))
    return _blocks_to_html(blocks)


def _classify_paragraph(paragraph, paragraph_html: str) -> dict[str, Any]:
    level = paragraph.level if paragraph.level is not None else 0
    for prefix in _UL_PREFIXES:
        if paragraph.text.startswith(prefix):
            return {"kind": "ul", "level": level, "inner_html": _strip_prefix_html(paragraph_html, prefix)}
    if _OL_PREFIX_RE.match(paragraph.text):
        match = _OL_PREFIX_RE.match(paragraph.text)
        assert match is not None
        return {
            "kind": "ol",
            "level": level,
            "inner_html": _strip_prefix_html(paragraph_html, match.group(0)),
        }
    heading_tag = _detect_heading(paragraph)
    if heading_tag:
        return {"kind": heading_tag, "inner_html": paragraph_html}
    return {"kind": "p", "inner_html": paragraph_html}


def _blocks_to_html(blocks: list[dict[str, Any]]) -> str:
    html_parts: list[str] = []
    list_stack: list[dict[str, Any]] = []

    def close_to_level(target_level: int) -> None:
        while list_stack and list_stack[-1]["level"] >= target_level:
            html_parts.append(f"</{list_stack.pop()['kind']}>")

    for block in blocks:
        kind = block["kind"]
        if kind in {"ul", "ol"}:
            level = block["level"]
            while list_stack and (
                list_stack[-1]["level"] > level
                or (list_stack[-1]["level"] == level and list_stack[-1]["kind"] != kind)
            ):
                html_parts.append(f"</{list_stack.pop()['kind']}>")
            if not list_stack or list_stack[-1]["level"] < level or list_stack[-1]["kind"] != kind:
                html_parts.append(f"<{kind}>")
                list_stack.append({"kind": kind, "level": level})
            html_parts.append(f"<li>{block['inner_html']}</li>")
            continue

        close_to_level(-1)
        if kind in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            html_parts.append(f"<{kind}>{block['inner_html']}</{kind}>")
        else:
            html_parts.append(f"<p>{block['inner_html']}</p>")

    close_to_level(-1)
    return "".join(html_parts)


def _detect_heading(paragraph) -> str | None:
    if not paragraph.runs:
        return None
    first = paragraph.runs[0]
    if first.font.bold is not True:
        return None
    size = first.font.size
    if size is None:
        return None
    pt = float(size.pt)
    if pt >= 30:
        return "h1"
    if pt >= 26:
        return "h2"
    if pt >= 22:
        return "h3"
    if pt >= 18:
        return "h4"
    if pt >= 15:
        return "h5"
    if pt >= 12:
        return "h6"
    return None


def _paragraph_html(paragraph) -> str:
    fragments: list[str] = []
    run_iter = iter(paragraph.runs)
    for child in paragraph._p:
        local = child.tag.split("}")[-1]
        if local == "r":
            run = next(run_iter, None)
            if run is not None:
                fragments.append(_run_html(run))
        elif local == "br":
            fragments.append("<br/>")
    return "".join(fragments)


def _run_html(run) -> str:
    text = html.escape(run.text or "")
    if text == "":
        return ""

    css_props: list[str] = []
    if run.font.name:
        css_props.append(f"font-family:{run.font.name}")
    if run.font.size:
        css_props.append(f"font-size:{round(float(run.font.size.pt), 2)}pt")
    color = _run_color(run)
    if color:
        css_props.append(f"color:#{color}")
    if css_props:
        text = f'<span style="{"; ".join(css_props)}">{text}</span>'

    if run.font.name == "Courier New":
        text = f"<code>{text}</code>"

    if run.font.underline is True:
        text = f"<u>{text}</u>"
    if run.font.italic is True:
        text = f"<i>{text}</i>"
    if run.font.bold is True:
        text = f"<b>{text}</b>"

    rpr = run.font._element
    if rpr.get("strike") == "sngStrike":
        text = f"<s>{text}</s>"

    baseline = rpr.get("baseline")
    if baseline == "30000":
        text = f"<sup>{text}</sup>"
    elif baseline == "-25000":
        text = f"<sub>{text}</sub>"

    if run.hyperlink.address:
        text = f'<a href="{html.escape(run.hyperlink.address, quote=True)}">{text}</a>'

    return text


def _run_color(run) -> str | None:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    return str(rgb)


def _strip_prefix_html(content: str, prefix: str) -> str:
    escaped = html.escape(prefix)
    if content.startswith(escaped):
        return content[len(escaped) :]
    return content.replace(escaped, "", 1)
