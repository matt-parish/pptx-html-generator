"""HTML to python-pptx rich text parser."""

from __future__ import annotations

from typing import Any

from bs4 import BeautifulSoup
from bs4.element import NavigableString, Tag
from pptx.util import Pt

from .styles import parse_css_color, parse_css_font_size, parse_css_inline_style

_BOLD_TAGS = {"b", "strong"}
_ITALIC_TAGS = {"i", "em"}
_UNDERLINE_TAGS = {"u"}
_STRIKE_TAGS = {"s", "del", "strike"}
_SUPERSCRIPT_BASELINE = "30000"
_SUBSCRIPT_BASELINE = "-25000"
_HEADING_TAG_TO_PT = {
    "h1": 32,
    "h2": 28,
    "h3": 24,
    "h4": 20,
    "h5": 16,
    "h6": 13,
}
_UL_BULLETS = ["\u2022 ", "\u2013 ", "\u25e6 "]


def render_html_to_text_frame(text_frame, content: str, base_styles: dict[str, Any] | None = None) -> None:
    """Render HTML content into a text frame."""
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]

    if not content:
        return

    root_styles = _default_styles()
    if base_styles:
        root_styles.update(base_styles)

    # Fast path: no obvious markup, treat as plain text.
    if "<" not in content and "&" not in content:
        if _styles_are_empty(root_styles):
            paragraph.text = content
        else:
            _add_run(paragraph, content, root_styles)
        return

    soup = BeautifulSoup(content, "lxml")
    container = soup.body if soup.body is not None else soup
    _process_children(text_frame, container, paragraph, root_styles, [], [])


def _default_styles() -> dict[str, Any]:
    return {
        "bold": None,
        "italic": None,
        "underline": None,
        "strikethrough": None,
        "font_name": None,
        "font_size": None,
        "font_color": None,
        "hyperlink": None,
        "baseline": None,
    }


def _process_children(
    text_frame,
    parent,
    paragraph,
    styles: dict[str, Any],
    list_stack: list[str],
    ol_counters: list[int],
):
    after_block = False

    for node in parent.children:
        if isinstance(node, NavigableString):
            text = str(node)
            if text == "":
                continue
            if after_block:
                paragraph = text_frame.add_paragraph()
                after_block = False
            _add_run(paragraph, text, styles)
            continue

        if not isinstance(node, Tag):
            continue

        tag_name = node.name.lower()
        tag_styles = _merge_tag_styles(styles, node, tag_name)

        if tag_name == "br":
            if after_block:
                paragraph = text_frame.add_paragraph()
                after_block = False
            paragraph.add_line_break()
            continue

        if tag_name in {"ul", "ol"}:
            list_stack.append(tag_name)
            if tag_name == "ol":
                ol_counters.append(0)
            paragraph = _process_children(
                text_frame, node, paragraph, tag_styles, list_stack, ol_counters
            )
            if tag_name == "ol":
                ol_counters.pop()
            list_stack.pop()
            after_block = True
            continue

        if tag_name == "li":
            if after_block or _paragraph_has_content(paragraph):
                paragraph = text_frame.add_paragraph()
            _configure_list_paragraph(paragraph, list_stack)
            _add_run(paragraph, _list_prefix(list_stack, ol_counters), styles)
            paragraph = _process_children(
                text_frame, node, paragraph, tag_styles, list_stack, ol_counters
            )
            after_block = True
            continue

        if tag_name in _HEADING_TAG_TO_PT:
            if after_block or _paragraph_has_content(paragraph):
                paragraph = text_frame.add_paragraph()
            heading_styles = tag_styles.copy()
            heading_styles["bold"] = True
            if heading_styles["font_size"] is None:
                heading_styles["font_size"] = Pt(_HEADING_TAG_TO_PT[tag_name])
            paragraph = _process_children(
                text_frame, node, paragraph, heading_styles, list_stack, ol_counters
            )
            after_block = True
            continue

        if tag_name == "p":
            if after_block or _paragraph_has_content(paragraph):
                paragraph = text_frame.add_paragraph()
            paragraph = _process_children(
                text_frame, node, paragraph, tag_styles, list_stack, ol_counters
            )
            after_block = True
            continue

        if after_block:
            paragraph = text_frame.add_paragraph()
            after_block = False
        paragraph = _process_children(text_frame, node, paragraph, tag_styles, list_stack, ol_counters)

    return paragraph


def _merge_tag_styles(styles: dict[str, Any], node: Tag, tag_name: str) -> dict[str, Any]:
    merged = styles.copy()
    if tag_name in _BOLD_TAGS:
        merged["bold"] = True
    if tag_name in _ITALIC_TAGS:
        merged["italic"] = True
    if tag_name in _UNDERLINE_TAGS:
        merged["underline"] = True
    if tag_name in _STRIKE_TAGS:
        merged["strikethrough"] = True
    if tag_name == "code":
        merged["font_name"] = "Courier New"
    if tag_name == "a":
        href = node.attrs.get("href")
        if isinstance(href, str) and href.strip():
            merged["hyperlink"] = href.strip()
    if tag_name == "sup":
        merged["baseline"] = _SUPERSCRIPT_BASELINE
    if tag_name == "sub":
        merged["baseline"] = _SUBSCRIPT_BASELINE

    style_value = node.attrs.get("style")
    if isinstance(style_value, str):
        _apply_css_overrides(merged, style_value)
    return merged


def _apply_css_overrides(merged: dict[str, Any], style_value: str) -> None:
    declarations = parse_css_inline_style(style_value)

    font_weight = declarations.get("font-weight")
    if font_weight:
        normalized = font_weight.strip().lower()
        if normalized in {"normal", "400"}:
            merged["bold"] = False
        elif normalized in {"bold", "700", "800", "900"}:
            merged["bold"] = True

    font_style = declarations.get("font-style")
    if font_style:
        normalized = font_style.strip().lower()
        if normalized == "normal":
            merged["italic"] = False
        elif normalized == "italic":
            merged["italic"] = True

    text_decoration = declarations.get("text-decoration")
    if text_decoration:
        normalized = text_decoration.strip().lower()
        if "none" in normalized:
            merged["underline"] = False
            merged["strikethrough"] = False
        if "underline" in normalized:
            merged["underline"] = True
        if "line-through" in normalized:
            merged["strikethrough"] = True

    font_family = declarations.get("font-family")
    if font_family:
        merged["font_name"] = font_family.split(",")[0].strip().strip("\"'")

    color_value = declarations.get("color")
    if color_value:
        try:
            merged["font_color"] = parse_css_color(color_value)
        except ValueError:
            pass

    size_value = declarations.get("font-size")
    if size_value:
        try:
            merged["font_size"] = parse_css_font_size(size_value)
        except ValueError:
            pass


def _add_run(paragraph, text: str, styles: dict[str, Any]) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.bold = styles["bold"]
    run.font.italic = styles["italic"]
    run.font.underline = styles["underline"]

    if styles["font_name"] is not None:
        run.font.name = styles["font_name"]
    if styles["font_size"] is not None:
        run.font.size = styles["font_size"]
    if styles["font_color"] is not None:
        run.font.color.rgb = styles["font_color"]
    if styles["hyperlink"] is not None:
        run.hyperlink.address = styles["hyperlink"]

    if styles["strikethrough"] is True:
        run.font._element.set("strike", "sngStrike")
    elif styles["strikethrough"] is False:
        run.font._element.set("strike", "noStrike")

    if styles["baseline"] is not None:
        run.font._element.set("baseline", styles["baseline"])


def _styles_are_empty(styles: dict[str, Any]) -> bool:
    return all(value is None for value in styles.values())


def _configure_list_paragraph(paragraph, list_stack: list[str]) -> None:
    if not list_stack:
        return
    paragraph.level = max(0, len(list_stack) - 1)


def _list_prefix(list_stack: list[str], ol_counters: list[int]) -> str:
    if not list_stack:
        return ""
    current = list_stack[-1]
    depth = max(0, len(list_stack) - 1)
    if current == "ol":
        ol_index = len(ol_counters) - 1
        if ol_index >= 0:
            ol_counters[ol_index] += 1
            return f"{ol_counters[ol_index]}. "
        return "1. "
    return _UL_BULLETS[depth % len(_UL_BULLETS)]


def _paragraph_has_content(paragraph) -> bool:
    if paragraph.text:
        return True
    return any(run.text for run in paragraph.runs)
