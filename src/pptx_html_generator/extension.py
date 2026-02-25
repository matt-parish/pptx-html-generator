"""Monkey-patch extension: adds an ``html`` property to python-pptx TextFrame.

Call :func:`install_html_extension` once (or rely on the package ``__init__``
which calls it automatically) to attach the property.  After that:

    shape.text_frame.html = "<b>Hello</b> world"   # setter: render HTML in
    html = shape.text_frame.html                    # getter: extract HTML out

Call :func:`uninstall_html_extension` to remove the patch (primarily useful in
tests or if you need to swap implementations at runtime).
"""

from __future__ import annotations

_SENTINEL = "pptx_html_generator"


def install_html_extension() -> None:
    """Attach a ``.html`` property to ``pptx.text.text.TextFrame``.

    Safe to call multiple times — subsequent calls are no-ops when the
    property was already installed by this package.

    Raises ``RuntimeError`` if ``TextFrame.html`` already exists and was
    *not* installed by this package, to avoid silently clobbering third-party
    extensions.
    """
    from pptx.text.text import TextFrame

    if hasattr(TextFrame, "html"):
        existing_sentinel = getattr(TextFrame, "_html_installed_by", None)
        if existing_sentinel == _SENTINEL:
            return  # idempotent — we installed it, nothing to do
        raise RuntimeError(
            f"TextFrame.html already exists (installed by: {existing_sentinel!r}). "
            "Cannot install pptx_html_generator extension without risking a conflict."
        )

    from .html_parser import render_html_to_text_frame
    from .reverse import _text_frame_to_html

    def _html_getter(self) -> str:
        return _text_frame_to_html(self)

    def _html_setter(self, value: str) -> None:
        render_html_to_text_frame(self, value)

    TextFrame.html = property(_html_getter, _html_setter)
    TextFrame._html_installed_by = _SENTINEL


def uninstall_html_extension() -> None:
    """Remove the ``.html`` property from ``TextFrame`` if we installed it.

    This is a no-op if the extension is not currently installed by this package.
    Primarily useful for test isolation.
    """
    from pptx.text.text import TextFrame

    if getattr(TextFrame, "_html_installed_by", None) != _SENTINEL:
        return

    del TextFrame.html
    del TextFrame._html_installed_by
