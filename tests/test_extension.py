"""Tests for the TextFrame.html monkey-patch extension."""

import pytest
from pptx import Presentation
from pptx.text.text import TextFrame

import pptx_html_generator  # auto-installs the extension
from pptx_html_generator.extension import install_html_extension, uninstall_html_extension


def _text_frame():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 3000000, 1000000)
    return shape.text_frame


# ---------------------------------------------------------------------------
# Installation state
# ---------------------------------------------------------------------------


def test_property_is_present_after_import():
    assert hasattr(TextFrame, "html")
    assert isinstance(TextFrame.__dict__["html"], property)


def test_sentinel_attribute_is_set():
    assert TextFrame._html_installed_by == "pptx_html_generator"


def test_install_is_idempotent():
    # Calling install again must not raise and must leave the property intact.
    install_html_extension()
    install_html_extension()
    assert hasattr(TextFrame, "html")


def test_uninstall_removes_property():
    uninstall_html_extension()
    try:
        assert not hasattr(TextFrame, "html")
        assert not hasattr(TextFrame, "_html_installed_by")
    finally:
        # Always reinstall so the rest of the suite is unaffected.
        install_html_extension()


def test_uninstall_is_idempotent():
    uninstall_html_extension()
    uninstall_html_extension()  # second call must be a no-op, not raise
    install_html_extension()  # restore


def test_install_raises_on_foreign_html_attribute():
    """If something else already put .html on TextFrame before us, we should refuse."""
    # Start from a clean (uninstalled) state, then simulate a foreign occupant.
    uninstall_html_extension()
    TextFrame.html = "foreign"
    try:
        with pytest.raises(RuntimeError, match="already exists"):
            install_html_extension()
    finally:
        del TextFrame.html
        # Restore so the rest of the suite is unaffected.
        install_html_extension()


# ---------------------------------------------------------------------------
# Setter behaviour
# ---------------------------------------------------------------------------


def test_setter_renders_bold():
    tf = _text_frame()
    tf.html = "<b>Bold text</b>"
    runs = tf.paragraphs[0].runs
    bold_run = next(r for r in runs if r.text == "Bold text")
    assert bold_run.font.bold is True


def test_setter_renders_italic():
    tf = _text_frame()
    tf.html = "<i>Italic</i>"
    assert tf.paragraphs[0].runs[0].font.italic is True


def test_setter_renders_multiple_paragraphs():
    tf = _text_frame()
    tf.html = "<p>First</p><p>Second</p>"
    assert len(tf.paragraphs) == 2
    assert tf.paragraphs[0].text == "First"
    assert tf.paragraphs[1].text == "Second"


def test_setter_clears_previous_content():
    tf = _text_frame()
    tf.html = "<p>Old</p>"
    tf.html = "<p>New</p>"
    texts = [p.text for p in tf.paragraphs]
    assert "New" in texts
    assert "Old" not in texts


def test_setter_accepts_plain_text():
    tf = _text_frame()
    tf.html = "Plain text no markup"
    assert tf.paragraphs[0].text == "Plain text no markup"


def test_setter_accepts_empty_string():
    tf = _text_frame()
    tf.html = ""
    # clear() leaves one empty paragraph — that is expected behaviour
    assert len(tf.paragraphs) == 1
    assert tf.paragraphs[0].text == ""


# ---------------------------------------------------------------------------
# Getter behaviour
# ---------------------------------------------------------------------------


def test_getter_returns_string():
    tf = _text_frame()
    tf.html = "<p>Hello</p>"
    result = tf.html
    assert isinstance(result, str)


def test_getter_contains_text_content():
    tf = _text_frame()
    tf.html = "<p>Hello world</p>"
    assert "Hello world" in tf.html


def test_getter_preserves_bold_tag():
    tf = _text_frame()
    tf.html = "<p><b>Bold</b></p>"
    assert "<b>" in tf.html
    assert "Bold" in tf.html


# ---------------------------------------------------------------------------
# Roundtrip
# ---------------------------------------------------------------------------


def test_roundtrip_simple_paragraph():
    tf = _text_frame()
    tf.html = "<p>Simple</p>"
    assert "<p>Simple</p>" == tf.html


def test_roundtrip_bold_inline():
    tf = _text_frame()
    tf.html = "<p>Hello <b>World</b></p>"
    result = tf.html
    assert "Hello" in result
    assert "<b>" in result
    assert "World" in result


def test_roundtrip_unordered_list():
    tf = _text_frame()
    tf.html = "<ul><li>One</li><li>Two</li></ul>"
    result = tf.html
    assert "<ul>" in result
    assert "<li>" in result
    assert "One" in result
    assert "Two" in result
