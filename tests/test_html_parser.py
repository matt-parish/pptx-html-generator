from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx_html_generator.html_parser import render_html_to_text_frame


def _text_frame():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 3000000, 1000000)
    return shape.text_frame


def test_plain_text_fast_path():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "Revenue grew 12%")

    assert len(text_frame.paragraphs) == 1
    assert text_frame.paragraphs[0].text == "Revenue grew 12%"


def test_inline_formatting_runs():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "A <b>bold</b> and <i>italic</i> and <u>u</u>")
    runs = text_frame.paragraphs[0].runs

    assert [run.text for run in runs] == ["A ", "bold", " and ", "italic", " and ", "u"]
    assert runs[1].font.bold is True
    assert runs[3].font.italic is True
    assert runs[5].font.underline is True


def test_br_inserts_soft_return():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "Line one<br>Line two")

    assert len(text_frame.paragraphs) == 1
    xml = text_frame.paragraphs[0]._p.xml
    assert "<a:br/>" in xml
    assert "Line one" in text_frame.paragraphs[0].text
    assert "Line two" in text_frame.paragraphs[0].text


def test_p_creates_new_paragraphs():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "<p>First</p><p>Second</p>")

    assert len(text_frame.paragraphs) == 2
    assert text_frame.paragraphs[0].text == "First"
    assert text_frame.paragraphs[1].text == "Second"


def test_whitespace_between_inline_tags_is_preserved():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "<b>hello</b> <i>world</i>")

    assert text_frame.paragraphs[0].text == "hello world"


def test_css_color_and_font_size_and_family():
    text_frame = _text_frame()
    render_html_to_text_frame(
        text_frame,
        '<span style="color:#00ff00; font-size: 20pt; font-family: Arial">styled</span>',
    )
    run = text_frame.paragraphs[0].runs[0]

    assert run.text == "styled"
    assert run.font.color.rgb == RGBColor(0x00, 0xFF, 0x00)
    assert run.font.size.pt == 20
    assert run.font.name == "Arial"


def test_css_override_disables_bold():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, '<b style="font-weight: normal">text</b>')
    run = text_frame.paragraphs[0].runs[0]
    assert run.font.bold is False


def test_hyperlink_code_strike_and_baseline():
    text_frame = _text_frame()
    render_html_to_text_frame(
        text_frame,
        '<a href="https://example.com"><code>code</code></a> <s>x</s> <sup>2</sup> <sub>n</sub>',
    )
    runs = text_frame.paragraphs[0].runs

    assert runs[0].text == "code"
    assert runs[0].hyperlink.address == "https://example.com"
    assert runs[0].font.name == "Courier New"

    strike_run = next(run for run in runs if run.text == "x")
    assert 'strike="sngStrike"' in strike_run.font._element.xml

    sup_run = next(run for run in runs if run.text == "2")
    sub_run = next(run for run in runs if run.text == "n")
    assert 'baseline="30000"' in sup_run.font._element.xml
    assert 'baseline="-25000"' in sub_run.font._element.xml


def test_lists_and_headings():
    text_frame = _text_frame()
    render_html_to_text_frame(
        text_frame,
        "<h2>Heading</h2><ul><li>One</li><li>Two<ul><li>Nested</li></ul></li></ul><ol><li>First</li><li>Second</li></ol>",
    )

    assert text_frame.paragraphs[0].runs[0].text == "Heading"
    assert text_frame.paragraphs[0].runs[0].font.bold is True
    assert round(text_frame.paragraphs[0].runs[0].font.size.pt) == 28

    assert text_frame.paragraphs[1].text.startswith("\u2022 One")
    assert text_frame.paragraphs[2].text.startswith("\u2022 Two")
    assert text_frame.paragraphs[3].text.startswith("\u2013 Nested")
    assert text_frame.paragraphs[3].level == 1
    assert text_frame.paragraphs[4].text.startswith("1. First")
    assert text_frame.paragraphs[5].text.startswith("2. Second")


def test_nested_inline_styles_with_css_override():
    text_frame = _text_frame()
    render_html_to_text_frame(
        text_frame,
        '<b>Bold <i>Both <span style="font-style:normal">BoldOnly</span></i> End</b>',
    )
    runs = text_frame.paragraphs[0].runs

    assert [run.text for run in runs] == ["Bold ", "Both ", "BoldOnly", " End"]
    assert runs[0].font.bold is True
    assert runs[1].font.bold is True and runs[1].font.italic is True
    assert runs[2].font.bold is True and runs[2].font.italic is False
    assert runs[3].font.bold is True and runs[3].font.italic is None


def test_malformed_overlapping_tags_gracefully_degrade():
    text_frame = _text_frame()
    render_html_to_text_frame(text_frame, "<p><b>bold <i>mix</b> tail</i> done</p>")
    runs = text_frame.paragraphs[0].runs

    assert text_frame.paragraphs[0].text == "bold mix tail done"
    assert runs[0].text == "bold " and runs[0].font.bold is True
    assert runs[1].text == "mix" and runs[1].font.bold is True and runs[1].font.italic is True
    assert runs[2].text == " tail done"


def test_deeply_nested_mixed_lists_keep_numbering_and_levels():
    text_frame = _text_frame()
    render_html_to_text_frame(
        text_frame,
        "<ol><li>A<ul><li><b>B</b><ol><li>C</li></ol></li></ul></li><li>D</li></ol>",
    )

    assert text_frame.paragraphs[0].text.startswith("1. A")
    assert text_frame.paragraphs[1].text.startswith("\u2013 B")
    assert text_frame.paragraphs[1].runs[1].font.bold is True
    assert text_frame.paragraphs[2].text.startswith("1. C")
    assert text_frame.paragraphs[2].level == 2
    assert text_frame.paragraphs[3].text.startswith("2. D")
