from pathlib import Path

from pptx import Presentation

from pptx_html_generator.generator import generate_pptx


def test_generate_pptx_plain_text(tmp_path: Path):
    spec = {
        "presentation": {"width": "10in", "height": "7.5in"},
        "slides": [
            {
                "layout": "blank",
                "elements": [
                    {
                        "type": "textbox",
                        "position": {
                            "left": "1in",
                            "top": "1in",
                            "width": "5in",
                            "height": "2in",
                        },
                        "content": "Revenue grew 12%",
                    }
                ],
            }
        ],
    }
    output = tmp_path / "out.pptx"
    generate_pptx(spec, output)

    assert output.exists()

    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes[0].text == "Revenue grew 12%"


def test_generate_pptx_rich_text_tier2(tmp_path: Path):
    spec = {
        "presentation": {"width": "10in", "height": "7.5in"},
        "slides": [
            {
                "layout": "blank",
                "elements": [
                    {
                        "type": "textbox",
                        "position": {
                            "left": "1in",
                            "top": "1in",
                            "width": "8in",
                            "height": "2in",
                        },
                        "content": '<p><span style="color:#008000">green</span> <a href="https://example.com">link</a></p><p>A<br/>B</p>',
                    }
                ],
            }
        ],
    }
    output = tmp_path / "out-rich.pptx"
    generate_pptx(spec, output)

    prs = Presentation(str(output))
    tf = prs.slides[0].shapes[0].text_frame

    assert len(tf.paragraphs) == 2
    assert "green" in tf.paragraphs[0].text
    assert "link" in tf.paragraphs[0].text
    assert tf.paragraphs[0].runs[2].hyperlink.address == "https://example.com"
    assert "<a:br/>" in tf.paragraphs[1]._p.xml


def test_generate_pptx_with_defaults_and_style(tmp_path: Path):
    spec = {
        "presentation": {
            "width": "10in",
            "height": "7.5in",
            "defaults": {"font_name": "Arial", "font_size": "18pt", "font_color": "#008000"},
        },
        "slides": [
            {
                "layout": "blank",
                "elements": [
                    {
                        "type": "textbox",
                        "position": {
                            "left": "1in",
                            "top": "1in",
                            "width": "8in",
                            "height": "2in",
                        },
                        "style": {"alignment": "center", "vertical_anchor": "middle", "word_wrap": True},
                        "content": "Styled plain text",
                    }
                ],
            }
        ],
    }
    output = tmp_path / "out-style.pptx"
    generate_pptx(spec, output)

    prs = Presentation(str(output))
    tf = prs.slides[0].shapes[0].text_frame
    run = tf.paragraphs[0].runs[0]

    assert run.font.name == "Arial"
    assert round(run.font.size.pt) == 18
    assert str(run.font.color.rgb) == "008000"
    assert tf.paragraphs[0].alignment is not None
